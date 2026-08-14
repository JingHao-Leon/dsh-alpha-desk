#!/usr/bin/env python3
"""Convert the Alpha Desk PPTD deck into a native, fully editable PPTX.

Scoped to the element/feature subset this deck uses:
text (rich <p>/<span>/<strong>/<a>), shape (rect/ellipse/homePlate),
line (sharp multi-segment, optional end arrow), image (cover crop), table.

Usage: python3 pptd_to_pptx.py <deck.pptd> <output.pptx>
Requires: python-pptx, PyYAML.
"""
from __future__ import annotations

import html.parser
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu

# ---------------------------------------------------------------- theme


def load_deck(manifest: Path):
    deck = yaml.safe_load(manifest.read_text(encoding="utf-8"))
    base = manifest.parent
    pages = [yaml.safe_load((base / p).read_text(encoding="utf-8")) for p in deck["pages"]]
    return deck, pages


class Theme:
    def __init__(self, deck):
        t = deck.get("theme") or {}
        self.colors = t.get("colors") or {}
        self.text_styles = t.get("textStyles") or {}

    def color(self, value: str | None, bg: str = "#FFFFFF") -> RGBColor | None:
        if not value:
            return None
        if value.startswith("$"):
            value = self.colors.get(value[1:], "#000000")
        value = value.strip()
        if re.fullmatch(r"#[0-9a-fA-F]{8}", value):  # HEX8 -> blend over bg
            fg, alpha = value[1:7], int(value[7:9], 16) / 255.0
            br, bg_, bb = int(bg[1:3], 16), int(bg[3:5], 16), int(bg[5:7], 16)
            fr, fg_, fb = int(fg[0:2], 16), int(fg[2:4], 16), int(fg[4:6], 16)
            r = round(fr * alpha + br * (1 - alpha))
            g = round(fg_ * alpha + bg_ * (1 - alpha))
            b = round(fb * alpha + bb * (1 - alpha))
            return RGBColor(r, g, b)
        return RGBColor.from_string(value.lstrip("#"))


# ---------------------------------------------------------------- rich text


class RichTextParser(html.parser.HTMLParser):
    """Parse our small HTML subset into [(text, style_dict, href), ...] paragraphs."""

    def __init__(self):
        super().__init__()
        self.paragraphs: list[dict] = []
        self._p: dict | None = None
        self._span_style: dict = {}
        self._bold = False
        self._href = None

    @staticmethod
    def _parse_style(attr: str | None) -> dict:
        out = {}
        if not attr:
            return out
        for part in attr.split(";"):
            if ":" in part:
                k, v = part.split(":", 1)
                out[k.strip()] = v.strip()
        return out

    def handle_starttag(self, tag, attrs):
        attrs = dict(attrs)
        if tag == "p":
            self._p = {"style": self._parse_style(attrs.get("style")), "runs": []}
        elif tag == "span":
            self._span_style = self._parse_style(attrs.get("style"))
        elif tag == "strong":
            self._bold = True
        elif tag == "a":
            self._href = attrs.get("href")
        elif tag == "br":
            self._add_run("\n")

    def handle_endtag(self, tag):
        if tag == "p" and self._p is not None:
            self.paragraphs.append(self._p)
            self._p = None
        elif tag == "span":
            self._span_style = {}
        elif tag == "strong":
            self._bold = False
        elif tag == "a":
            self._href = None

    def _add_run(self, data):
        if self._p is None:
            self._p = {"style": {}, "runs": []}
        self._p["runs"].append(
            {"text": data, "span": dict(self._span_style), "bold": self._bold, "href": self._href}
        )

    def handle_data(self, data):
        if data.strip() or data == "\n":
            self._add_run(data)


def parse_rich(text: str) -> list[dict]:
    if "<" not in text:  # plain text: one paragraph per line
        return [{"style": {}, "runs": [{"text": line, "span": {}, "bold": False, "href": None}]}
                for line in text.split("\n") if line.strip()]
    parser = RichTextParser()
    parser.feed(text)
    return parser.paragraphs


# ---------------------------------------------------------------- helpers

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
         "justify": PP_ALIGN.JUSTIFY}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
SHAPES = {"rect": MSO_SHAPE.RECTANGLE, "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
          "ellipse": MSO_SHAPE.OVAL, "homePlate": MSO_SHAPE.PENTAGON}


def set_font(run, *, size, color, bold, italic, family, spacing=None):
    f = run.font
    f.size = Pt(size)
    if color is not None:
        f.color.rgb = color
    f.bold = bold
    f.italic = italic
    if family:
        f.name = family  # sets a:latin in schema position
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            latin = rPr.find(qn("a:latin"))
            if latin is not None:
                latin.addnext(ea)  # a:ea must immediately follow a:latin
            else:
                rPr.append(ea)
        ea.set("typeface", family)
    if spacing:
        run._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))


def apply_fill(shape, fill, theme: Theme, bg: str):
    if not fill or fill.get("type") != "solid":
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.color(fill.get("color"), bg)


def apply_border(shape, border, theme: Theme, bg: str):
    ln = shape.line
    if not border:
        ln.fill.background()
        return
    ln.color.rgb = theme.color(border.get("color", "#000000"), bg)
    ln.width = Pt(border.get("width", 1))
    if border.get("style") == "dash":
        ln_el = shape._element.spPr.find(qn("a:ln"))
        dash = ln_el.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln_el.append(dash)


def style_of(content: dict, theme: Theme) -> dict:
    """Resolve a text content dict against its referenced theme text style."""
    base = {}
    ref = content.get("style")
    if ref and ref.startswith("$"):
        base = dict(theme.text_styles.get(ref[1:], {}))
    merged = {**base, **{k: v for k, v in content.items() if v is not None}}
    return merged


def add_text(slide, el, theme: Theme, bg: str):
    content = el["content"]
    st = style_of(content, theme)
    x, y, w, h = el["bounds"]
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = box.text_frame
    tf.word_wrap = content.get("wrap", True)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    halign, valign = content.get("align", ["left", "top"])
    tf.vertical_anchor = ANCHOR.get(valign, MSO_ANCHOR.TOP)

    base_size = st.get("fontSize", 18)
    base_color = theme.color(st.get("color", "#000000"), bg)
    base_bold = st.get("bold", False)
    base_family = st.get("fontFamily")
    base_spacing = st.get("letterSpacing")
    base_line = st.get("lineHeight")

    for i, para in enumerate(parse_rich(content["text"])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pstyle = para["style"]
        p.alignment = ALIGN.get(pstyle.get("text-align", halign), PP_ALIGN.LEFT)
        lh = pstyle.get("line-height", base_line)
        if lh:
            p.line_spacing = float(lh)
        if pstyle.get("margin-top"):
            p.space_before = Pt(float(pstyle["margin-top"].removesuffix("px")))
        for r in para["runs"]:
            run = p.add_run()
            run.text = r["text"]
            span = r["span"]
            size = span.get("font-size")
            color = span.get("color")
            set_font(
                run,
                size=float(size.removesuffix("px")) if size else base_size,
                color=theme.color(color, bg) if color else base_color,
                bold=True if r["bold"] else base_bold,
                italic=False,
                family=span.get("font-family", base_family),
                spacing=base_spacing,
            )
            if r.get("href"):
                # after fonts/colors: keeps rPr child order schema-valid,
                # and the explicit solidFill wins over the theme hyperlink color
                run.hyperlink.address = r["href"]
    return box


def add_shape(slide, el, theme: Theme, bg: str):
    x, y, w, h = el["bounds"]
    kind = SHAPES.get(el["shapeName"], MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
    apply_fill(shape, el.get("fill"), theme, bg)
    apply_border(shape, el.get("border"), theme, bg)
    shape.shadow.inherit = False
    return shape


def add_line(slide, el, theme: Theme, bg: str):
    x, y, w, h = el["bounds"]
    vw, vh = el["viewBox"]
    pts = []
    for pair in el["points"].split():
        px, py = pair.split(",")
        pts.append((x + float(px) / vw * w, y + float(py) / vh * h))
    border = el.get("border") or {}
    prev = None
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
        conn.line.color.rgb = theme.color(border.get("color", "#000000"), bg)
        conn.line.width = Pt(border.get("width", 1))
        conn.shadow.inherit = False
        ln = conn._element.spPr.find(qn("a:ln"))
        if border.get("style") == "dash":
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        if i == len(pts) - 2 and (el.get("arrow") or [None, None])[1]:
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "arrow", "w": "med", "len": "med"}))
        prev = conn
    return prev


def add_image(slide, el, manifest_dir: Path, theme: Theme, bg: str):
    x, y, w, h = el["bounds"]
    src = manifest_dir / el["src"]
    crop = el.get("crop") or {}
    pic = slide.shapes.add_picture(str(src), Pt(x), Pt(y), Pt(w), Pt(h))
    # apply proportional crop first (fractions of source), then cover-fit
    pic.crop_left = crop.get("left", 0)
    pic.crop_right = crop.get("right", 0)
    pic.crop_top = crop.get("top", 0)
    pic.crop_bottom = crop.get("bottom", 0)
    if (el.get("fit") or {}).get("mode") == "cover":
        from PIL import Image
        iw, ih = Image.open(src).size
        iw *= 1 - pic.crop_left - pic.crop_right
        ih *= 1 - pic.crop_top - pic.crop_bottom
        src_ratio, dst_ratio = iw / ih, w / h
        if src_ratio > dst_ratio:  # too wide -> crop sides further
            keep = dst_ratio / src_ratio
            extra = (1 - keep) / 2
            pic.crop_left = pic.crop_left + extra * (1 - pic.crop_left - pic.crop_right)
            pic.crop_right = pic.crop_right + extra * (1 - pic.crop_left - pic.crop_right)
        else:  # too tall -> crop top/bottom further
            keep = src_ratio / dst_ratio
            extra = (1 - keep) / 2
            pic.crop_top = pic.crop_top + extra * (1 - pic.crop_top - pic.crop_bottom)
            pic.crop_bottom = pic.crop_bottom + extra * (1 - pic.crop_top - pic.crop_bottom)
    border = el.get("border")
    if border:
        pic.line.color.rgb = theme.color(border.get("color", "#D8DEE4"), bg)
        pic.line.width = Pt(border.get("width", 1))
    pic.shadow.inherit = False
    return pic


def set_cell_border(cell, edge, color="D8DEE4", width_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    tag = {"left": "a:lnL", "right": "a:lnR", "top": "a:lnT", "bottom": "a:lnB"}[edge]
    for old in tcPr.findall(qn(tag)):
        tcPr.remove(old)
    ln = tcPr.makeelement(qn(tag), {"w": str(int(width_pt * 12700)), "cap": "flat"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr)
    ln.append(fill)
    tcPr.append(ln)


def add_table(slide, el, theme: Theme, bg: str):
    x, y, w, h = el["bounds"]
    rows, cols = len(el["rows"]), len(el["columnWidths"])
    gfx = slide.shapes.add_table(rows, cols, Pt(x), Pt(y), Pt(w), Pt(h))
    table = gfx.table
    for i, ratio in enumerate(el["columnWidths"]):
        table.columns[i].width = Emu(int(Pt(w) * ratio))
    for i, ratio in enumerate(el["rowHeights"]):
        table.rows[i].height = Emu(int(Pt(h) * ratio))
    for r, row in enumerate(el["rows"]):
        for c, spec in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Pt(6)
            cell.margin_top = cell.margin_bottom = Pt(2)
            cell.vertical_anchor = ANCHOR.get((spec.get("align") or ["center", "middle"])[1],
                                              MSO_ANCHOR.MIDDLE)
            if spec.get("fill") and spec["fill"].get("type") == "solid":
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.color(spec["fill"].get("color"), bg)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for edge in ("left", "right", "top", "bottom"):
                set_cell_border(cell, edge)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = ALIGN.get((spec.get("align") or ["center", "middle"])[0], PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = spec.get("text", "")
            set_font(run,
                     size=spec.get("fontSize", 12),
                     color=theme.color(spec.get("color", "#33404F"), bg),
                     bold=spec.get("bold", False),
                     italic=False,
                     family=spec.get("fontFamily"))
    return gfx


# ---------------------------------------------------------------- main


def convert(manifest: Path, output: Path):
    deck, pages = load_deck(manifest)
    theme = Theme(deck)
    width, height = deck.get("size", [960, 540])
    prs = Presentation()
    prs.slide_width = Pt(width)
    prs.slide_height = Pt(height)
    blank = prs.slide_layouts[6]

    for page in pages:
        slide = prs.slides.add_slide(blank)
        bg_fill = page.get("background") or {"type": "solid", "color": "#FFFFFF"}
        bg_hex = bg_fill.get("color", "#FFFFFF")
        if bg_hex.startswith("$"):
            bg_hex = theme.colors.get(bg_hex[1:], "#FFFFFF")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#")[:6])

        for el in page.get("elements", []):
            kind = el["elementType"]
            if kind == "text":
                add_text(slide, el, theme, bg_hex)
            elif kind == "shape":
                add_shape(slide, el, theme, bg_hex)
            elif kind == "line":
                add_line(slide, el, theme, bg_hex)
            elif kind == "image":
                add_image(slide, el, manifest.parent, theme, bg_hex)
            elif kind == "table":
                add_table(slide, el, theme, bg_hex)
            else:
                print(f"  skip unsupported element: {kind}", file=sys.stderr)

    prs.save(output)
    print(f"saved {output} ({len(pages)} slides)")


if __name__ == "__main__":
    convert(Path(sys.argv[1]), Path(sys.argv[2]))
